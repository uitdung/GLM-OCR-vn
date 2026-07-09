#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Inspect position/size/font của từng placeholder trong các layout cần dùng."""
from pptx import Presentation
from pptx.util import Emu

SRC = r"c:\uit\nlp\GLM-OCR\tools\slides\Buổi 1-1.pptx"
prs = Presentation(SRC)

LAYOUTS_OF_INTEREST = ["Tiêu đề", "Noi dung", "So sanh 1", "So sanh 2", "Noi dung 2",
                      "Hình ảnh", "1_Mục lục", "Tiêu đề chương"]

def emu_in(v): return round(Emu(v).inches, 2) if v is not None else None

for lay in prs.slide_layouts:
    if lay.name not in LAYOUTS_OF_INTEREST:
        continue
    print(f"\n══ LAYOUT '{lay.name}' ══")
    for ph in lay.placeholders:
        try:
            l = emu_in(ph.left); t = emu_in(ph.top)
            w = emu_in(ph.width); h = emu_in(ph.height)
            fmt = ph.placeholder_format
            # font info from first run if any
            font_info = ""
            try:
                if ph.has_text_frame:
                    for p in ph.text_frame.paragraphs:
                        for r in p.runs:
                            sz = r.font.size.pt if r.font.size else None
                            bold = r.font.bold
                            nm = r.font.name
                            col = None
                            try:
                                col = str(r.font.color.rgb) if r.font.color and r.font.color.type else None
                            except: pass
                            font_info = f"sz={sz} bold={bold} font={nm} color={col}"
                            break
                        if font_info: break
            except: pass
            print(f"  idx={fmt.idx} type={fmt.type} name='{ph.name}'")
            print(f"    pos=({l}, {t}) size=({w}×{h})  {font_info}")
        except Exception as e:
            print(f"  idx={ph.placeholder_format.idx} ERROR {e}")

# Còn inspect luôn slide thực tế để xem cách fill text ra sao
print("\n══ SAMPLE FILLED SLIDE (slide 3 'Noi dung') ══")
s = prs.slides[2]
for sh in s.shapes:
    if not sh.has_text_frame: continue
    txt = sh.text_frame.text.strip().replace("\n", " | ")
    if not txt: continue
    print(f"  shape '{sh.name}' pos=({emu_in(sh.left)},{emu_in(sh.top)}) size=({emu_in(sh.width)}×{emu_in(sh.height)})")
    print(f"    text: {txt[:100]}")
