#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Phân tích theme + layout của Buổi 1-1.pptx để copy sang deck mới."""
import sys, zipfile, re
from pptx import Presentation
from pptx.util import Emu

SRC = r"c:\uit\nlp\GLM-OCR\tools\slides\Buổi 1-1.pptx"

prs = Presentation(SRC)
print(f"═ Slide size: {Emu(prs.slide_width).inches:.2f} × {Emu(prs.slide_height).inches:.2f} inches")
print(f"═ Total slides: {len(prs.slides)}")
print(f"═ Slide masters: {len(prs.slide_masters)}")
print(f"═ Slide layouts: {len(prs.slide_layouts)}")

print("\n── LAYOUTS ──")
for i, lay in enumerate(prs.slide_layouts):
    phs = [(p.placeholder_format.idx, p.placeholder_format.type) for p in lay.placeholders]
    print(f"  [{i}] '{lay.name}'  placeholders={phs}")

print("\n── MASTERS ──")
for i, m in enumerate(prs.slide_masters):
    print(f"  master[{i}]: '{m.name}'  bg_shapes={len(m.shapes)}")

# Theme colors
print("\n── THEME COLORS (master[0]) ──")
try:
    master = prs.slide_masters[0]
    from pptx.oxml.ns import qn
    theme = master.element.getroottree()
    # Try reading theme XML directly
    import os
    z = zipfile.ZipFile(SRC)
    theme_files = [n for n in z.namelist() if n.startswith("ppt/theme/")]
    print(f"  theme files: {theme_files}")
    for tf in theme_files[:1]:
        data = z.read(tf).decode("utf-8")
        # extract color scheme
        colors = re.findall(r'<a:srgbClr val="([0-9A-Fa-f]{6})"/>', data)
        print(f"  {tf}: {len(colors)} srgbClr refs")
        # Clarify color scheme elements
        scheme = re.findall(r'<a:(dk1|lt1|dk2|lt2|accent[1-6]|hlink|folHlink)>.*?val="([0-9A-Fa-f]{6})"', data, re.S)
        for name, val in scheme[:12]:
            print(f"    {name} = #{val}")
        # fonts
        fonts = re.findall(r'typeface="([^"]+)"', data)
        uniq = []
        for f in fonts:
            if f not in uniq: uniq.append(f)
        print(f"  typefaces (unique): {uniq[:20]}")
except Exception as e:
    print(f"  theme parse error: {e}")

print("\n── PER-SLIDE OVERVIEW ──")
for idx, slide in enumerate(prs.slides):
    layout_name = slide.slide_layout.name
    n_shapes = len(slide.shapes)
    n_pics = len([s for s in slide.shapes if s.shape_type == 13])
    n_txt = len([s for s in slide.shapes if s.has_text_frame])
    # First text snippet
    first_text = ""
    for s in slide.shapes:
        if s.has_text_frame and s.text_frame.text.strip():
            first_text = s.text_frame.text.strip().replace("\n", " | ")[:80]
            break
    print(f"  slide {idx+1}: layout='{layout_name}' shapes={n_shapes} txt={n_txt} pics={n_pics} | {first_text}")
