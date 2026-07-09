#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Sinh deck báo cáo tiến độ giữa kỳ CHO GLM-OCR-vn — DÙNG TEMPLATE UIT.
Phiên bản cô đọng: text ngắn, chèn ảnh minh hoạ, gộp tổng kết+kế hoạch.

9 slide:
  1  Bìa                       layout 'Tiêu đề'
  2  Bài toán OCR tiếng Việt   layout 'Noi dung'
  3  Các mô hình OCR + ảnh     layout 'Noi dung 2'  (vlm_ocr_evolution.png)
  4  Kiến trúc GLM tổng quan   layout 'Noi dung 2'  (glm_ocr_architecture.png)
  5  ViT vs Projector          layout 'So sanh 1'
  6  LLM Decoder + LoRA        layout 'Noi dung 2'  (lora_architecture.png)
  7  Tiến trình train + ảnh    layout 'Noi dung 2'  (finetune_pipeline.png)
  8  Tổng kết + Kế hoạch (gộp) layout 'So sanh 1'
  9  Cảm ơn                    layout 'Noi dung'
"""
import os, shutil, tempfile, zipfile, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Paths ─────────────────────────────────────────────────────────
ROOT = r"c:\uit\nlp\GLM-OCR"
SRC  = os.path.join(ROOT, "tools", "slides", "Buổi 1-1.pptx")
OUT  = os.path.join(ROOT, "tools", "slides", "bao_cao_giua_ky.pptx")
IMG  = {
    "evolution": os.path.join(ROOT, "vlm_ocr_evolution.png"),
    "architecture": os.path.join(ROOT, "glm_ocr_architecture.png"),
    "lora": os.path.join(ROOT, "lora_architecture.png"),
    "pipeline": os.path.join(ROOT, "finetune_pipeline.png"),
}
for k, v in IMG.items():
    assert os.path.exists(v), f"Missing image: {v}"

# ── Màu UIT ───────────────────────────────────────────────────────
NAVY   = RGBColor(0x1C, 0x30, 0x5E)
DK     = RGBColor(0x2A, 0x2F, 0x4F)
BLUE   = RGBColor(0x00, 0x71, 0xFF)
ORANGE = RGBColor(0xFA, 0xAB, 0x78)
RED    = RGBColor(0xEC, 0x71, 0x71)
GREEN  = RGBColor(0x41, 0x85, 0x5B)
PURPLE = RGBColor(0x47, 0x00, 0xD8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x60, 0x60, 0x70)
LBLUE  = RGBColor(0xE8, 0xF1, 0xFF)


# ── Bước 1: tạo template UIT sạch (xoá slide, giữ master/layouts) ─
def strip_slides_from_template(src_pptx, dst_pptx):
    tmp = tempfile.mkdtemp(prefix="pptx_clean_")
    try:
        with zipfile.ZipFile(src_pptx, "r") as z:
            z.extractall(tmp)
        slides_dir = os.path.join(tmp, "ppt", "slides")
        rels_dir   = os.path.join(slides_dir, "_rels")
        if os.path.isdir(slides_dir):
            for fn in os.listdir(slides_dir):
                if re.match(r"slide\d+\.xml$", fn):
                    os.remove(os.path.join(slides_dir, fn))
        if os.path.isdir(rels_dir):
            for fn in os.listdir(rels_dir):
                if re.match(r"slide\d+\.xml\.rels$", fn):
                    os.remove(os.path.join(rels_dir, fn))
        # presentation.xml: rỗng sldIdLst
        with open(os.path.join(tmp, "ppt", "presentation.xml"), "r", encoding="utf-8") as f:
            c = f.read()
        c = re.sub(r"<p:sldIdLst>.*?</p:sldIdLst>", "<p:sldIdLst></p:sldIdLst>", c, flags=re.S)
        with open(os.path.join(tmp, "ppt", "presentation.xml"), "w", encoding="utf-8") as f:
            f.write(c)
        # presentation.xml.rels: bỏ slide rels
        with open(os.path.join(tmp, "ppt", "_rels", "presentation.xml.rels"), "r", encoding="utf-8") as f:
            rc = f.read()
        rc = re.sub(r'<Relationship[^>]*Target="slides/slide\d+\.xml"[^/]*/>', "", rc)
        with open(os.path.join(tmp, "ppt", "_rels", "presentation.xml.rels"), "w", encoding="utf-8") as f:
            f.write(rc)
        # [Content_Types].xml
        with open(os.path.join(tmp, "[Content_Types].xml"), "r", encoding="utf-8") as f:
            cc = f.read()
        cc = re.sub(r'<Override PartName="/ppt/slides/slide\d+\.xml"[^/]*/>', "", cc)
        with open(os.path.join(tmp, "[Content_Types].xml"), "w", encoding="utf-8") as f:
            f.write(cc)
        if os.path.exists(dst_pptx):
            os.remove(dst_pptx)
        with zipfile.ZipFile(dst_pptx, "w", zipfile.ZIP_DEFLATED) as z:
            ct = os.path.join(tmp, "[Content_Types].xml")
            z.write(ct, "[Content_Types].xml")
            for r, _, files in os.walk(tmp):
                for fn in files:
                    full = os.path.join(r, fn)
                    arc = os.path.relpath(full, tmp).replace(os.sep, "/")
                    if arc == "[Content_Types].xml":
                        continue
                    z.write(full, arc)
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


CLEAN = os.path.join(tempfile.gettempdir(), "uit_template_clean.pptx")
strip_slides_from_template(SRC, CLEAN)
print(f"✓ Template UIT sạch → {CLEAN}")

prs = Presentation(CLEAN)
SW, SH = prs.slide_width, prs.slide_height


def lay(name):
    for l in prs.slide_layouts:
        if l.name == name:
            return l
    raise ValueError(name)

L_COVER = lay("Tiêu đề")
L_BODY  = lay("Noi dung")
L_BODY2 = lay("Noi dung 2")     # title trái + content trái + OBJECT phải (ảnh)
L_CMP1  = lay("So sanh 1")      # 2 cột content


# ── Helper ────────────────────────────────────────────────────────
def norm(item, sz0, col0, bld0, aln0):
    """item: str | dict | tuple(text,[size],[color],[bold],[indent],[align])."""
    txt, sz, col, bld, ind, aln = "", sz0, col0, bld0, 0, aln0
    if isinstance(item, str):
        txt = item
    elif isinstance(item, dict):
        txt = item.get("text", ""); sz = item.get("size", sz0)
        col = item.get("color", col0); bld = item.get("bold", bld0)
        ind = item.get("indent", 0); aln = item.get("align", aln0)
    elif isinstance(item, tuple):
        if len(item) > 0: txt = item[0]
        if len(item) > 1: sz = item[1]
        if len(item) > 2: col = item[2]
        if len(item) > 3: bld = item[3]
        if len(item) > 4: ind = item[4]
        if len(item) > 5: aln = item[5]
    return txt, sz, col, bld, ind, aln


def fill_ph(ph, lines, size=14, color=DK, bold=False, align=PP_ALIGN.LEFT,
            bullet=True, gap=4):
    tf = ph.text_frame; tf.clear(); tf.word_wrap = True
    for i, item in enumerate(lines):
        txt, sz, col, bld, ind, aln = norm(item, size, color, bold, align)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = aln
        p.space_after = Pt(gap); p.space_before = Pt(0)
        prefix = ""
        if bullet:
            prefix = "" if ind < 0 else ("   • " if ind >= 1 else "▸ ")
        r = p.add_run(); r.text = prefix + txt
        r.font.size = Pt(sz); r.font.bold = bld
        r.font.color.rgb = col; r.font.name = "Calibri Light"
    return ph


def set_title(slide, title_text, color=NAVY, size=None):
    """Set title (placeholder idx=0 hoặc TITLE) — bìa không có nên skip."""
    for ph in slide.placeholders:
        if ph.placeholder_format.type == 1 and ph.has_text_frame:
            tf = ph.text_frame; tf.text = title_text
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.name = "Calibri Light"; r.font.bold = True
                    try: r.font.color.rgb = color
                    except: pass
                    if size: r.font.size = Pt(size)
            return ph
    return None


def add_image_contain(slide, pic_path, x, y, w, h):
    """Chèn ảnh fit-contain trong khung (x,y,w,h), giữ aspect, căn giữa."""
    from PIL import Image
    iw, ih = Image.open(pic_path).size
    scale = min(w / iw, h / ih)
    nw = int(iw * scale); nh = int(ih * scale)
    nx = int(x + (w - nw) / 2); ny = int(y + (h - nh) / 2)
    return slide.shapes.add_picture(pic_path, nx, ny, width=nw, height=nh)


def get_ph(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — BÌA
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(L_COVER)
fill_ph(get_ph(s, 13), [("BÁO CÁO TIẾN ĐỘ GIỮA KỲ", 30, NAVY, True)],
        bullet=False)
fill_ph(get_ph(s, 16),
        [("Fine-tune mô hình GLM-OCR\ncho OCR Tiếng Việt", 24, BLUE, True)],
        bullet=False)
fill_ph(get_ph(s, 14),
        [("Khắc phục lỗi dấu thanh khi áp dụng VLM", 14, GRAY, False)],
        bullet=False)
fill_ph(get_ph(s, 15), [("tuandung-specominc", 14, NAVY, True)],
        align=PP_ALIGN.CENTER, bullet=False)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — BÀI TOÁN OCR TIẾNG VIỆT (Noi dung, cô đọng)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(L_BODY)
set_title(s, "Bài toán: OCR tiếng Việt")
fill_ph(get_ph(s, 1), [
    ("Đặc thù dấu thanh tiếng Việt", 20, NAVY, True, -1),
    ("6 dấu thanh + 7 nhóm nguyên âm kép (ă â ê ô ơ ư đ)", 16, DK, False, 0),
    ("Một ký tự = tổ hợp 3 lớp: chữ cái + dấu mũ + dấu thanh", 16, DK, False, 0),
    ("Dấu chỉ vài pixel → mờ/nhiễu dễ mất dấu", 16, DK, False, 0),
    ("", 8),
    ("Vấn đề với mô hình OCR gốc", 20, NAVY, True, -1),
    ("GLM-OCR pre-train cho tiếng Anh/Trung → yếu dấu Việt", 16, DK, False, 0),
    ("Mất dấu = thay đổi nghĩa ('than' ≠ 'thần' ≠ 'thăn')", 16, RED, False, 0),
    ("False Positive: model tự thêm dấu sai", 16, RED, False, 0),
], size=16, gap=8)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — CÁC MÔ HÌNH OCR + ảnh (Noi dung 2)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(L_BODY2)
set_title(s, "Các mô hình\nOCR & VLM", color=NAVY)
# placeholder idx 1 (left content 4.3×4.44)
fill_ph(get_ph(s, 1), [
    ("Tiến hóa 2017→2026", 16, NAVY, True, -1),
    ("CRNN→TrOCR→Donut→GLM-OCR", 14, DK, False, 0),
    ("Ba hướng tiếp cận:", 16, NAVY, True, -1),
    ("Pipeline (PP-OCRv6)", 14, BLUE, True, 0),
    ("VLM end-to-end (GLM-OCR) ★", 14, GREEN, True, 0),
    ("Chọn GLM-OCR:", 16, ORANGE, True, -1),
    ("Tận dụng LLM, linh hoạt, phù hợp VN", 13, DK, False, 0),
], size=14, gap=5)
# placeholder idx 15 (right OBJECT 6.85×6.53) — chèn ảnh
right_ph = get_ph(s, 15)
# Xoá placeholder rỗng rồi vẽ ảnh trong khung đó để vừa slide
rp_x, rp_y = Emu(right_ph.left), Emu(right_ph.top)
rp_w, rp_h = Emu(right_ph.width), Emu(right_ph.height)
rp_x_in, rp_y_in = rp_x / 914400, rp_y / 914400
rp_w_in, rp_h_in = rp_w / 914400, rp_h / 914400
add_image_contain(s, IMG["evolution"],
                  Inches(rp_x_in), Inches(rp_y_in),
                  Inches(rp_w_in), Inches(rp_h_in))


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — KIẾN TRÚC TỔNG QUAN + ảnh
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(L_BODY2)
set_title(s, "Kiến trúc\nGLM-OCR (1/2)", color=NAVY)
fill_ph(get_ph(s, 1), [
    ("VLM end-to-end ~1.1 tỷ TS", 16, NAVY, True, -1),
    ("① ViT CogViT — 24L, h=1024", 14, BLUE, True, 0),
    ("image 336 · patch 14 · merge 2", 12, DK, False, 1),
    ("② Projector — bridge ViT↔LLM", 14, NAVY, True, 0),
    ("d_vision → d_llm = 1536", 12, DK, False, 1),
    ("③ LLM Decoder GLM-0.5B", 14, GREEN, True, 0),
    ("16L · GQA · vocab 59392", 12, DK, False, 1),
    ("MTP loss — học phụ thuộc dài", 12, DK, False, 1),
], size=14, gap=5)
right_ph = get_ph(s, 15)
add_image_contain(s, IMG["architecture"],
                  Inches(Emu(right_ph.left)/914400),
                  Inches(Emu(right_ph.top)/914400),
                  Inches(Emu(right_ph.width)/914400),
                  Inches(Emu(right_ph.height)/914400))


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — ViT vs PROJECTOR (So sanh 1, 2 cột, text cô đọng)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(L_CMP1)
set_title(s, "Kiến trúc (2/2) — ViT & Projector")
# idx 1 left content, idx 15 right content (5.69×5.41 each)
fill_ph(get_ph(s, 1), [
    ("① Vision Encoder — CogViT", 18, NAVY, True, -1),
    ("24 self-attention, hidden=1024", 15, DK, False, 0),
    ("image=336, patch=14", 15, DK, False, 0),
    ("Nén → ~999 token / ảnh 768×1024", 15, DK, False, 0),
    ("out_hidden=1536 (khớp LLM)", 15, DK, False, 0),
    ("", 8),
    ("⚠ Vì sao dấu thanh khó?", 15, ORANGE, True, -1),
    ("Dấu mũ/râu chỉ vài pixel", 14, RED, False, 0),
    ("→ dễ bị nén mất khi merge", 14, RED, False, 0),
], size=15, gap=6)
fill_ph(get_ph(s, 15), [
    ("② Projector — Bridge thị giác", 18, NAVY, True, -1),
    ("Chuyển đặc trưng ViT → embedding LLM", 15, DK, False, 0),
    ("Vị trí xây 'từ điển thị giác' VN", 15, DK, False, 0),
    ("", 8),
    ("2 chiến lược finetune", 15, BLUE, True, -1),
    ("Stage 1: mở projector (học dấu)", 14, DK, False, 1),
    ("Stage 2: đóng băng (ổn định)", 14, DK, False, 1),
    ("", 8),
    ("Cấu hình v3 đang chạy", 15, ORANGE, True, -1),
    ("FREEZE cả ViT + Projector", 14, RED, False, 1),
    ("→ chỉ train LLM bằng LoRA", 14, GREEN, True, 1),
], size=15, gap=6)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — LLM DECODER + LoRA  (Noi dung 2 + ảnh lora)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(L_BODY2)
set_title(s, "LLM Decoder\n& LoRA", color=NAVY)
fill_ph(get_ph(s, 1), [
    ("③ LLM — GLM-0.5B", 16, NAVY, True, -1),
    ("16L · GQA 16Q/8KV", 14, DK, False, 0),
    ("vocab 59392 · mRoPE", 14, DK, False, 0),
    ("", 8),
    ("Quyết định finetune v3", 16, ORANGE, True, -1),
    ("ViT — FROZEN", 14, DK, False, 0),
    ("Projector — FROZEN", 14, DK, False, 0),
    ("LLM + LoRA — TRAINED", 14, GREEN, True, 0),
    ("rank 8 · alpha 32", 12, GRAY, False, 1),
], size=14, gap=5)
right_ph = get_ph(s, 15)
add_image_contain(s, IMG["lora"],
                  Inches(Emu(right_ph.left)/914400),
                  Inches(Emu(right_ph.top)/914400),
                  Inches(Emu(right_ph.width)/914400),
                  Inches(Emu(right_ph.height)/914400))


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — TIẾN TRÌNH TRAIN + ảnh pipeline
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(L_BODY2)
set_title(s, "Tiến trình\ntrain", color=NAVY)
fill_ph(get_ph(s, 1), [
    ("Pipeline (notebook)", 16, NAVY, True, -1),
    ("GPU → data → cài LF →", 13, DK, False, 0),
    ("train → eval → dashboard", 13, DK, False, 0),
    ("", 8),
    ("Cấu hình LoRA v3", 16, ORANGE, True, -1),
    ("rank 8, α32, target='all'", 13, DK, False, 0),
    ("freeze ViT + projector", 13, RED, False, 0),
    ("lr 1e-5 · cutoff 1024", 13, DK, False, 0),
    ("Eff.BS 16 · grad_ckpt=true", 13, DK, False, 0),
    ("", 8),
    ("Đánh giá", 16, GREEN, True, -1),
    ("CER · WER · EM% · DA%", 13, BLUE, True, 0),
    ("⏳ Đang thu checkpoint", 12, ORANGE, False, 0),
], size=13, gap=4)
right_ph = get_ph(s, 15)
add_image_contain(s, IMG["pipeline"],
                  Inches(Emu(right_ph.left)/914400),
                  Inches(Emu(right_ph.top)/914400),
                  Inches(Emu(right_ph.width)/914400),
                  Inches(Emu(right_ph.height)/914400))


# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — TỔNG KẾT + KẾ HOẠCH (GỘP, So sanh 1)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(L_CMP1)
set_title(s, "Tổng kết & kế hoạch")
# Cột trái: Tiến độ (3 trạng thái cô đọng)
fill_ph(get_ph(s, 1), [
    ("✓ ĐÃ HOÀN THÀNH", 18, GREEN, True, -1),
    ("Khảo sát + sơ đồ kiến trúc", 14, DK, False, 0),
    ("Pipeline data ShareGPT 3 nguồn", 14, DK, False, 0),
    ("YAML v3 + eval + dashboard", 14, DK, False, 0),
    ("", 8),
    ("⏳ ĐANG LÀM", 18, ORANGE, True, -1),
    ("Train 3 epoch × 3000 mẫu", 14, DK, False, 0),
    ("Thu checkpoint + so base vs FT", 14, DK, False, 0),
    ("", 8),
    ("🎯 KẾ HOẠCH (+Tuần)", 18, BLUE, True, -1),
    ("1-2: hoàn train · 3-4: mở projector", 13, DK, False, 0),
    ("5: dp chéo · 6: deploy + report", 13, DK, False, 0),
], size=14, gap=4)
# Cột phải: KPI + rủi ro
fill_ph(get_ph(s, 15), [
    ("Rủi ro & xử lý", 18, RED, True, -1),
    ("OOM T4 → bật grad_ckpt + fp16", 13, DK, False, 0),
    ("Freeze ViT hạn chế dấu → mở S2", 13, DK, False, 0),
    ("Overfit → theo dõi loss + CER", 13, DK, False, 0),
    ("", 8),
    ("KPI cuối kỳ mong đợi", 18, NAVY, True, -1),
    ("DA% ≥ 95% trên test set", 14, GREEN, True, 0),
    ("CER cải thiện rõ vs base", 14, GREEN, True, 0),
    ("Demo ảnh mới ngoài dataset", 14, GREEN, True, 0),
    ("", 8),
    ("📊 Tiến độ: ~55% — đúng kế hoạch", 16, ORANGE, True, -1),
], size=14, gap=4)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — CẢM ƠN (k co đọng)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(L_BODY)
set_title(s, "")
ph = get_ph(s, 1)
tf = ph.text_frame; tf.clear()
for i, (txt, sz, col, bld) in enumerate([
    ("Cảm ơn thầy/cô đã lắng nghe!", 36, NAVY, True),
    ("", 12, DK, False),
    ("Q & A", 28, BLUE, True),
    ("Sẵn sàng trả lời câu hỏi", 16, GRAY, False),
    ("", 14, DK, False),
    ("Dự án: uitdung/GLM-OCR-vn", 16, NAVY, True),
    ("Liên hệ: tuandung-specominc", 14, GRAY, False),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(6)
    r = p.add_run(); r.text = txt
    r.font.size = Pt(sz); r.font.bold = bld
    r.font.color.rgb = col; r.font.name = "Calibri Light"


prs.save(OUT)
print(f"\n✓ Saved: {OUT}")
print(f"  {len(prs.slides)} slide · Template UIT · cô đọng + ảnh")
